"""
PPTX Builder 테스트

DocumentOutput → PPTX 변환 엔진의 동작을 검증한다.
"""

import os
import tempfile

import pytest

from pptx import Presentation

# 프로젝트 모듈
from react_system.document_schema import (
    DocumentElement,
    DocumentOutput,
    ListContent,
    SectionOutput,
    TableColumn,
    TableContent,
    TextContent,
)
from react_system.tools.pptx_builder import build_pptx


# ─── 헬퍼 ───


def _minimal_doc(title="테스트 문서", doc_type="보고서", sections=None):
    """최소한의 DocumentOutput을 생성한다."""
    if sections is None:
        sections = []
    return DocumentOutput(title=title, doc_type=doc_type, sections=sections)


def _text_element(text, etype="paragraph", bold=False, font_size=None, alignment=None):
    """텍스트 기반 DocumentElement를 생성한다."""
    return DocumentElement(
        type=etype,
        content=TextContent(
            text=text, bold=bold, font_size=font_size, alignment=alignment
        ),
    )


def _table_element(columns, rows, caption=None, merge=None):
    """테이블 DocumentElement를 생성한다."""
    cols = [TableColumn(name=c) for c in columns]
    return DocumentElement(
        type="table",
        content=TableContent(columns=cols, rows=rows, caption=caption, merge=merge),
    )


def _list_element(items, list_type="bullet"):
    """목록 DocumentElement를 생성한다."""
    return DocumentElement(
        type="list",
        content=ListContent(items=items, list_type=list_type),
    )


# ─── 테스트 1: 타이틀 슬라이드 생성 ───


@pytest.mark.asyncio
async def test_should_create_title_slide_with_document_title_and_doc_type():
    """빈 섹션의 문서라도 타이틀 슬라이드가 생성되어야 한다."""
    doc = _minimal_doc(title="해양 엔진 부품 보고서", doc_type="보고서")

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = os.path.join(tmpdir, "test.pptx")
        result = await build_pptx(doc, output_path=output_path)

        assert result["status"] == "success"
        assert os.path.exists(output_path)

        prs = Presentation(output_path)
        # 최소 1개 슬라이드 (타이틀)
        assert len(prs.slides) >= 1

        title_slide = prs.slides[0]
        # 타이틀 슬라이드에 제목 텍스트가 있어야 한다
        title_shape = title_slide.shapes.title
        assert title_shape is not None
        assert "해양 엔진 부품 보고서" in title_shape.text


# ─── 테스트 2: 섹션 제목 슬라이드 생성 ───


@pytest.mark.asyncio
async def test_should_create_section_header_slide_for_each_section():
    """각 섹션마다 섹션 제목 슬라이드가 생성되어야 한다."""
    sections = [
        SectionOutput(section_id="s1", section_title="추진 배경", elements=[]),
        SectionOutput(section_id="s2", section_title="사업 개요", elements=[]),
    ]
    doc = _minimal_doc(title="사업 보고서", sections=sections)

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = os.path.join(tmpdir, "test.pptx")
        result = await build_pptx(doc, output_path=output_path)

        assert result["status"] == "success"

        prs = Presentation(output_path)
        # 타이틀(1) + 섹션 헤더(2) = 최소 3개 슬라이드
        assert len(prs.slides) >= 3

        # 두 번째 슬라이드 = 첫 번째 섹션 제목
        section1_slide = prs.slides[1]
        assert section1_slide.shapes.title is not None
        assert "추진 배경" in section1_slide.shapes.title.text

        # 세 번째 슬라이드 = 두 번째 섹션 제목
        section2_slide = prs.slides[2]
        assert section2_slide.shapes.title is not None
        assert "사업 개요" in section2_slide.shapes.title.text


# ─── 테스트 3: heading + paragraph 요소 → 본문 슬라이드 ───


@pytest.mark.asyncio
async def test_should_create_content_slide_with_heading_and_paragraph():
    """heading + paragraph 요소가 본문 슬라이드로 변환되어야 한다."""
    elements = [
        _text_element("추진 배경 상세", etype="heading"),
        _text_element("해양 엔진 부품 시장은 지속적으로 성장하고 있다.", etype="paragraph"),
    ]
    sections = [
        SectionOutput(section_id="s1", section_title="추진 배경", elements=elements),
    ]
    doc = _minimal_doc(title="보고서", sections=sections)

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = os.path.join(tmpdir, "test.pptx")
        result = await build_pptx(doc, output_path=output_path)

        assert result["status"] == "success"

        prs = Presentation(output_path)
        # 타이틀(1) + 섹션헤더(1) + 내용(1) = 최소 3개
        assert len(prs.slides) >= 3

        # 내용 슬라이드 (3번째)
        content_slide = prs.slides[2]
        # 슬라이드에서 heading이 제목으로 들어가야 한다
        assert content_slide.shapes.title is not None
        assert "추진 배경 상세" in content_slide.shapes.title.text

        # 본문 텍스트가 슬라이드 어딘가에 존재해야 한다
        all_text = " ".join(
            shape.text for shape in content_slide.shapes if shape.has_text_frame
        )
        assert "해양 엔진 부품 시장은 지속적으로 성장하고 있다" in all_text


# ─── 테스트 4: 테이블 요소 → 테이블 슬라이드 ───


@pytest.mark.asyncio
async def test_should_create_table_slide_with_header_and_data_rows():
    """table 요소가 테이블 슬라이드로 변환되어야 한다."""
    elements = [
        _text_element("부품 목록", etype="heading"),
        _table_element(
            columns=["부품명", "수량", "단가"],
            rows=[
                ["엔진 오일", "10", "50,000"],
                ["프로펠러", "2", "300,000"],
            ],
        ),
    ]
    sections = [
        SectionOutput(section_id="s1", section_title="부품 정보", elements=elements),
    ]
    doc = _minimal_doc(title="부품 보고서", sections=sections)

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = os.path.join(tmpdir, "test.pptx")
        result = await build_pptx(doc, output_path=output_path)

        assert result["status"] == "success"

        prs = Presentation(output_path)
        # 타이틀(1) + 섹션헤더(1) + 테이블(1) = 최소 3개
        assert len(prs.slides) >= 3

        # 테이블 슬라이드 찾기
        table_slide = prs.slides[2]
        # 테이블 shape가 있어야 한다
        table_shapes = [s for s in table_slide.shapes if s.has_table]
        assert len(table_shapes) >= 1

        table = table_shapes[0].table
        # 헤더(1) + 데이터(2) = 3행, 3열
        assert table.rows.__len__() == 3
        assert len(table.columns) == 3

        # 헤더 행 검증
        assert table.cell(0, 0).text == "부품명"
        assert table.cell(0, 1).text == "수량"
        assert table.cell(0, 2).text == "단가"

        # 데이터 행 검증
        assert table.cell(1, 0).text == "엔진 오일"
        assert table.cell(2, 1).text == "2"


# ─── 테스트 5: 목록(list) 요소 → 본문 슬라이드 ───


@pytest.mark.asyncio
async def test_should_create_content_slide_with_bullet_list():
    """bullet list 요소가 본문 슬라이드의 글머리표로 변환되어야 한다."""
    elements = [
        _text_element("주요 항목", etype="heading"),
        _list_element(["엔진 점검", "프로펠러 교체", "오일 교환"], list_type="bullet"),
    ]
    sections = [
        SectionOutput(section_id="s1", section_title="점검 항목", elements=elements),
    ]
    doc = _minimal_doc(title="점검 보고서", sections=sections)

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = os.path.join(tmpdir, "test.pptx")
        result = await build_pptx(doc, output_path=output_path)

        assert result["status"] == "success"

        prs = Presentation(output_path)
        # 타이틀(1) + 섹션헤더(1) + 내용(1) = 3개
        assert len(prs.slides) >= 3

        content_slide = prs.slides[2]
        all_text = " ".join(
            shape.text for shape in content_slide.shapes if shape.has_text_frame
        )
        # 글머리표 항목들이 텍스트에 포함되어야 한다
        assert "엔진 점검" in all_text
        assert "프로펠러 교체" in all_text
        assert "오일 교환" in all_text


@pytest.mark.asyncio
async def test_should_create_content_slide_with_numbered_list():
    """numbered list 요소가 번호 매김 목록으로 변환되어야 한다."""
    elements = [
        _text_element("절차", etype="heading"),
        _list_element(["전원 차단", "부품 분리", "교체 설치"], list_type="numbered"),
    ]
    sections = [
        SectionOutput(section_id="s1", section_title="절차", elements=elements),
    ]
    doc = _minimal_doc(title="절차서", sections=sections)

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = os.path.join(tmpdir, "test.pptx")
        result = await build_pptx(doc, output_path=output_path)

        assert result["status"] == "success"

        prs = Presentation(output_path)
        content_slide = prs.slides[2]
        all_text = " ".join(
            shape.text for shape in content_slide.shapes if shape.has_text_frame
        )
        # 번호가 붙은 항목이 포함되어야 한다
        assert "1. 전원 차단" in all_text
        assert "2. 부품 분리" in all_text
        assert "3. 교체 설치" in all_text


# ─── 테스트 6: 테이블 셀 병합 ───


@pytest.mark.asyncio
async def test_should_merge_table_cells_when_merge_spec_provided():
    """merge 정의가 있으면 셀 병합이 적용되어야 한다."""
    from react_system.document_schema import MergeCell

    elements = [
        _table_element(
            columns=["구분", "항목", "비고"],
            rows=[
                ["엔진", "오일 필터", "정기 교체"],
                ["엔진", "점화 플러그", "필요 시"],
                ["프로펠러", "블레이드", "점검"],
            ],
            merge=[
                MergeCell(start_row=1, start_col=0, end_row=2, end_col=0, value="엔진"),
            ],
        ),
    ]
    sections = [
        SectionOutput(section_id="s1", section_title="부품 표", elements=elements),
    ]
    doc = _minimal_doc(title="부품 보고서", sections=sections)

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = os.path.join(tmpdir, "test.pptx")
        result = await build_pptx(doc, output_path=output_path)

        assert result["status"] == "success"

        prs = Presentation(output_path)
        # 테이블 슬라이드 찾기 (타이틀 + 섹션헤더 + 테이블)
        table_slide = prs.slides[2]
        table_shapes = [s for s in table_slide.shapes if s.has_table]
        assert len(table_shapes) >= 1

        table = table_shapes[0].table
        # 병합된 셀의 텍스트 확인
        assert table.cell(1, 0).text == "엔진"


# ─── 테스트 7: tempfile 자동 생성 ───


@pytest.mark.asyncio
async def test_should_create_tempfile_when_no_output_path_given():
    """output_path를 지정하지 않으면 임시 파일로 생성되어야 한다."""
    doc = _minimal_doc(title="임시 파일 테스트")

    result = await build_pptx(doc, output_path=None)

    assert result["status"] == "success"
    assert result["file_path"] is not None
    assert result["file_path"].endswith(".pptx")
    assert os.path.exists(result["file_path"])

    # 정리
    os.unlink(result["file_path"])


# ─── 테스트 8: 에러 반환 ───


@pytest.mark.asyncio
async def test_should_return_error_status_on_failure():
    """잘못된 입력으로 에러가 발생하면 status: error를 반환해야 한다."""
    # DocumentOutput은 유효하지만 file_path를 쓸 수 없는 경로로 지정
    doc = _minimal_doc(title="에러 테스트")

    result = await build_pptx(doc, output_path="/nonexistent/dir/test.pptx")

    assert result["status"] == "error"
    assert "실패" in result["message"] or "오류" in result["message"]


# ─── 테스트 9: 복합 문서 (여러 섹션, 혼합 요소) ───


@pytest.mark.asyncio
async def test_should_handle_complex_document_with_mixed_elements():
    """여러 섹션과 다양한 요소 타입이 혼합된 문서를 처리해야 한다."""
    sections = [
        SectionOutput(
            section_id="s1",
            section_title="개요",
            elements=[
                _text_element("사업 개요", etype="heading"),
                _text_element("본 사업은 해양 엔진 부품을 대상으로 한다.", etype="paragraph"),
            ],
        ),
        SectionOutput(
            section_id="s2",
            section_title="부품 현황",
            elements=[
                _text_element("부품 목록", etype="heading"),
                _table_element(
                    columns=["품목", "가격"],
                    rows=[["엔진 오일", "50,000"], ["필터", "30,000"]],
                ),
                _text_element("주의사항", etype="heading"),
                _list_element(["정품만 사용", "유통기한 확인"], list_type="bullet"),
            ],
        ),
    ]
    doc = _minimal_doc(title="종합 보고서", doc_type="보고서", sections=sections)

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = os.path.join(tmpdir, "test.pptx")
        result = await build_pptx(doc, output_path=output_path)

        assert result["status"] == "success"

        prs = Presentation(output_path)
        # 타이틀(1) + 섹션1헤더(1) + 섹션1내용(1) + 섹션2헤더(1) + 테이블(1) + 주의사항(1) = 6개
        assert len(prs.slides) >= 6

        # 전체 슬라이드에서 핵심 텍스트 확인
        all_texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_texts.append(shape.text)
        combined = " ".join(all_texts)

        assert "종합 보고서" in combined
        assert "사업 개요" in combined
        assert "부품 목록" in combined
        assert "정품만 사용" in combined


# ─── 테스트 10: 테이블 캡션 ───


@pytest.mark.asyncio
async def test_should_display_table_caption_in_slide():
    """테이블에 caption이 있으면 슬라이드에 표시되어야 한다."""
    elements = [
        _table_element(
            columns=["품목", "단가"],
            rows=[["오일", "50,000"]],
            caption="[표 1] 부품 단가표",
        ),
    ]
    sections = [
        SectionOutput(section_id="s1", section_title="단가", elements=elements),
    ]
    doc = _minimal_doc(title="단가표", sections=sections)

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = os.path.join(tmpdir, "test.pptx")
        result = await build_pptx(doc, output_path=output_path)

        assert result["status"] == "success"

        prs = Presentation(output_path)
        table_slide = prs.slides[2]
        all_text = " ".join(
            shape.text for shape in table_slide.shapes if shape.has_text_frame
        )
        assert "[표 1] 부품 단가표" in all_text
