"""
PPTX 문서 빌드 엔진

DocumentOutput (JSON 스키마) → PPTX 파일 변환.
python-pptx 라이브러리를 사용하여 각 섹션을 슬라이드로 변환한다.
"""

import logging
import tempfile

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

from react_system.document_schema import (
    DocumentElement,
    DocumentOutput,
    SectionOutput,
    TableContent,
    TextContent,
)


# ─── 상수 ───

# 와이드스크린 슬라이드 크기 (13.333 x 7.5 인치)
SLIDE_WIDTH = Inches(13.333) if _PPTX_AVAILABLE else None
SLIDE_HEIGHT = Inches(7.5) if _PPTX_AVAILABLE else None

# 기본 폰트
DEFAULT_FONT = "Malgun Gothic"

# 헤더 배경색: 진한 파란색
HEADER_BG_COLOR = RGBColor(0x2F, 0x54, 0x96) if _PPTX_AVAILABLE else None
HEADER_FONT_COLOR = RGBColor(0xFF, 0xFF, 0xFF) if _PPTX_AVAILABLE else None


async def build_pptx(doc: DocumentOutput, output_path: str = None) -> dict:
    """
    DocumentOutput JSON → PPTX 파일 생성

    각 섹션을 1개 이상의 슬라이드로 변환:
    - 문서 제목 → 타이틀 슬라이드
    - 각 섹션 → 섹션 제목 슬라이드 + 내용 슬라이드들

    Returns:
        {"status": "success", "file_path": output_path, "message": "..."}
    """
    if not _PPTX_AVAILABLE:
        return {
            "status": "error",
            "file_path": None,
            "message": "python-pptx 라이브러리가 설치되지 않았습니다. "
                       "'pip install python-pptx'를 실행해주세요.",
        }

    try:
        prs = Presentation()

        # 와이드스크린 슬라이드 크기 설정
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        # 1) 타이틀 슬라이드
        _add_title_slide(prs, doc.title, doc.doc_type)

        # 2) 각 섹션 처리
        for section in doc.sections:
            _add_section_header_slide(prs, section.section_title)
            _add_section_content_slides(prs, section)

        # 파일 저장
        if output_path is None:
            tmp = tempfile.NamedTemporaryFile(
                suffix=".pptx", delete=False, prefix="doc_"
            )
            output_path = tmp.name
            tmp.close()

        prs.save(output_path)

        return {
            "status": "success",
            "file_path": output_path,
            "message": f"PPTX 파일이 생성되었습니다: {output_path}",
        }

    except Exception as e:
        logger.exception("PPTX 생성 중 오류 발생")
        return {
            "status": "error",
            "file_path": None,
            "message": f"PPTX 생성 실패: {str(e)}",
        }


# ─── 슬라이드 생성 함수 ───


def _add_title_slide(prs: "Presentation", title: str, subtitle: str) -> None:
    """문서 제목 슬라이드를 추가한다 (레이아웃 0)."""
    layout = prs.slide_layouts[0]  # 타이틀 슬라이드
    slide = prs.slides.add_slide(layout)

    # 제목
    if slide.shapes.title:
        slide.shapes.title.text = title
        _apply_font(slide.shapes.title.text_frame, font_size=Pt(36), bold=True)

    # 부제목 (doc_type)
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
        _apply_font(slide.placeholders[1].text_frame, font_size=Pt(20))


def _add_section_header_slide(prs: "Presentation", title: str) -> None:
    """섹션 제목 슬라이드를 추가한다 (레이아웃 2: Section Header)."""
    layout = prs.slide_layouts[2]  # 섹션 헤더
    slide = prs.slides.add_slide(layout)

    if slide.shapes.title:
        slide.shapes.title.text = title
        _apply_font(slide.shapes.title.text_frame, font_size=Pt(28), bold=True)


def _add_section_content_slides(prs: "Presentation", section: SectionOutput) -> None:
    """섹션의 요소들을 내용 슬라이드로 변환한다.

    heading 요소를 슬라이드 제목으로, 이후 paragraph/list 요소를 본문으로 그룹화한다.
    table 요소는 별도 슬라이드로 생성한다.
    """
    if not section.elements:
        return

    # 요소를 슬라이드 단위로 그룹화
    # heading이 나오면 새 슬라이드 그룹 시작
    slide_groups = _group_elements_into_slides(section.elements)

    for group in slide_groups:
        if group["type"] == "content":
            _add_content_slide(prs, group["heading"], group["body_elements"])
        elif group["type"] == "table":
            _add_table_slide(prs, group["heading"], group["table_content"])


def _group_elements_into_slides(
    elements: list[DocumentElement],
) -> list[dict]:
    """요소 리스트를 슬라이드 그룹으로 분리한다.

    Returns:
        [{"type": "content", "heading": str, "body_elements": [...]},
         {"type": "table", "heading": str, "table_content": TableContent}, ...]
    """
    groups = []
    current_heading = ""
    current_body = []

    for elem in elements:
        if elem.type == "heading":
            # 이전 그룹에 본문이 있으면 저장
            if current_body:
                groups.append({
                    "type": "content",
                    "heading": current_heading,
                    "body_elements": current_body,
                })
                current_body = []
            current_heading = elem.content.text

        elif elem.type == "table":
            # 이전 텍스트 그룹에 본문이 있으면 먼저 저장
            if current_body:
                groups.append({
                    "type": "content",
                    "heading": current_heading,
                    "body_elements": current_body,
                })
                current_body = []

            # 테이블 그룹 (heading을 테이블 제목으로 사용)
            groups.append({
                "type": "table",
                "heading": current_heading,
                "table_content": elem.content,
            })
            current_heading = ""

        else:
            # paragraph, list
            current_body.append(elem)

    # 남은 그룹 저장
    if current_heading or current_body:
        groups.append({
            "type": "content",
            "heading": current_heading,
            "body_elements": current_body,
        })

    return groups


def _add_content_slide(
    prs: "Presentation",
    heading: str,
    body_elements: list[DocumentElement],
) -> None:
    """heading + paragraph/list 요소 → 본문 슬라이드 (레이아웃 1: Title + Content)."""
    layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(layout)

    # 제목 설정
    if slide.shapes.title and heading:
        slide.shapes.title.text = heading
        _apply_font(slide.shapes.title.text_frame, font_size=Pt(24), bold=True)

    # 본문 내용 설정 (placeholder index 1 = 본문 영역)
    if len(slide.placeholders) > 1 and body_elements:
        body_placeholder = slide.placeholders[1]
        tf = body_placeholder.text_frame
        tf.clear()

        first = True
        for elem in body_elements:
            if elem.type == "paragraph":
                content = elem.content
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()

                p.text = content.text
                if content.bold:
                    for run in p.runs:
                        run.font.bold = True
                if content.font_size:
                    for run in p.runs:
                        run.font.size = Pt(content.font_size)
                if content.alignment:
                    p.alignment = _get_alignment(content.alignment)
                _apply_font_to_paragraph(p)

            elif elem.type == "list":
                content = elem.content
                for i, item in enumerate(content.items):
                    if first:
                        p = tf.paragraphs[0]
                        first = False
                    else:
                        p = tf.add_paragraph()

                    if content.list_type == "numbered":
                        p.text = f"{i + 1}. {item}"
                    else:
                        p.text = f"• {item}"
                    p.level = 1
                    _apply_font_to_paragraph(p)


def _add_table_slide(
    prs: "Presentation",
    heading: str,
    table_content: "TableContent",
) -> None:
    """테이블 요소 → 테이블 슬라이드.

    헤더 행(볼드, 진한 파란 배경) + 데이터 행으로 구성한다.
    """
    layout = prs.slide_layouts[5]  # Blank 레이아웃
    slide = prs.slides.add_slide(layout)

    # 테이블 제목 (heading 또는 caption이 있으면 텍스트 박스 추가)
    title_text = heading or ""
    if table_content.caption and not heading:
        title_text = table_content.caption
    elif table_content.caption and heading:
        title_text = f"{heading}\n{table_content.caption}"

    if title_text:
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(0.6)
        )
        tf = txBox.text_frame
        tf.text = title_text
        _apply_font(tf, font_size=Pt(24), bold=True)

    cols = table_content.columns
    rows_data = table_content.rows
    num_cols = len(cols)
    num_rows = len(rows_data) + 1  # +1 = 헤더 행

    # 테이블 위치 및 크기
    left = Inches(0.5)
    top = Inches(1.2) if title_text else Inches(0.5)
    width = Inches(12)
    height = Inches(0.4) * num_rows

    table_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, height
    )
    table = table_shape.table

    # 헤더 행 설정
    for col_idx, col_def in enumerate(cols):
        cell = table.cell(0, col_idx)
        cell.text = col_def.name
        # 헤더 스타일: 볼드, 흰 글씨, 진한 파란 배경
        _style_table_cell(cell, bold=True, font_color=HEADER_FONT_COLOR)
        _set_cell_bg_color(cell, HEADER_BG_COLOR)

    # 데이터 행 설정
    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            if col_idx < num_cols:
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = cell_text
                _style_table_cell(cell)

    # 셀 병합 처리
    if table_content.merge:
        for m in table_content.merge:
            start_cell = table.cell(m.start_row, m.start_col)
            end_cell = table.cell(m.end_row, m.end_col)
            start_cell.merge(end_cell)
            if m.value:
                start_cell.text = m.value


# ─── 유틸리티 ───


def _apply_font(text_frame, font_name=None, font_size=None, bold=None, color=None):
    """텍스트 프레임의 모든 run에 폰트 스타일을 적용한다."""
    font_name = font_name or DEFAULT_FONT
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            if font_size:
                run.font.size = font_size
            if bold is not None:
                run.font.bold = bold
            if color:
                run.font.color.rgb = color


def _apply_font_to_paragraph(paragraph, font_name=None, font_size=None):
    """개별 paragraph의 run에 기본 폰트를 적용한다."""
    font_name = font_name or DEFAULT_FONT
    for run in paragraph.runs:
        run.font.name = font_name
        if font_size:
            run.font.size = font_size


def _get_alignment(alignment: str):
    """문자열 정렬값을 PP_ALIGN 상수로 변환한다."""
    mapping = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }
    return mapping.get(alignment, PP_ALIGN.LEFT)


def _style_table_cell(cell, bold=None, font_color=None, font_size=None):
    """테이블 셀의 텍스트 스타일을 적용한다."""
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = DEFAULT_FONT
            run.font.size = font_size or Pt(11)
            if bold is not None:
                run.font.bold = bold
            if font_color:
                run.font.color.rgb = font_color


def _set_cell_bg_color(cell, color: "RGBColor"):
    """테이블 셀의 배경색을 설정한다.

    python-pptx는 셀 배경색 직접 설정 API가 없으므로
    XML을 직접 조작한다.
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # 기존 solidFill 제거
    for child in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(child)

    # 새 solidFill 추가
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", str(color))
